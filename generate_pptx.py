from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def save_presentation_from_data(summarized_data, save_path, images_path, template_path=None):
    """
    Генерирует файл презентации .pptx из словаря данных слайдов.

    :param summarized_data: Словарь данных слайдов после суммаризации.
    :param save_path: Путь для сохранения выходного файла .pptx.
    :param images_path: Путь к папке с изображениями.
    :param template_path: Путь к файлу шаблона .pptx, если он выбран.
    """
    # Если шаблон указан, используем его, иначе создаем новую презентацию
    if template_path and os.path.isfile(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    for slide_data in summarized_data:
        # Создаем слайд с титульной диаграммой
        slide_layout = prs.slide_layouts[1]  # Используем титульный слайд как основу
        slide = prs.slides.add_slide(slide_layout)

        # Добавляем заголовок
        title = slide.shapes.title
        title.text = slide_data['title']

        # Добавляем текстовые блоки в основной текстовый элемент слайда (placeholder[1])
        text_box = slide.shapes.placeholders[1]
        text_frame = text_box.text_frame
        text_frame.clear()  # Очищаем текстовый элемент, чтобы избежать дублирования текста

        # Добавляем суммаризированные параграфы
        for para in slide_data['paragraphs']:
            if para.strip():  # Проверяем, что параграф не пустой
                p = text_frame.add_paragraph()
                p.text = para
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.LEFT

        # Добавляем текст как есть
        for para in slide_data['text_as_is']:
            if para.strip():  # Проверяем, что текст не пустой
                p = text_frame.add_paragraph()
                p.text = para
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.LEFT

        # Добавляем изображение, если оно есть
        if slide_data['image']:
            image_path = os.path.join(images_path, slide_data['image'])
            if os.path.isfile(image_path):
                left = Inches(0.5)
                top = Inches(5)  # Ставим картинку в нижнюю часть слайда
                height = Inches(2.5)  # Высота изображения
                slide.shapes.add_picture(image_path, left, top, height=height)
            else:
                print(f"Изображение {slide_data['image']} не найдено в {images_path}")

    # Сохраняем презентацию
    prs.save(save_path)
    print(f"Презентация сохранена в {save_path}")
